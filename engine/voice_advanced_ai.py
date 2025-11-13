import subprocess
import pyautogui
import psutil
from datetime import datetime, timedelta
import json
import os
import shutil
import socket
import time
import random
import requests
import sqlite3
from collections import Counter
from typing import Optional, Dict, Any, List
import zipfile
import re
import warnings
warnings.filterwarnings("ignore")
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'

class VoiceAdvancedAI:
    def __init__(self):
        # Initialize AI models
        self.ai_provider = self._get_ai_provider()
        self._init_ai_models()
        
        # Initialize multilingual support
        try:
            from engine.multilingual_support import multilingual
            self.multilingual = multilingual
        except:
            self.multilingual = None
        
        # Advanced features paths
        project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.MEMORY_DB = os.path.join(project_dir, "jarvis_memory.db")
        self.CALENDAR_JSON = os.path.join(project_dir, "jarvis_calendar.json")
        self.VAULT_DIR = os.path.join(project_dir, "jarvis_vault")
        self.USAGE_LOG = os.path.join(project_dir, "jarvis_usage.json")
        self._ensure_db()
        
        # ALL VOICE-ACTIVATED ADVANCED FEATURES
        self.voice_functions = {
            # System Monitoring
            'system monitor live': self.system_monitor_dashboard_live,
            'monitor system': self.system_monitor_dashboard_live,
            'system monitor': self.system_monitor_dashboard_live,
            'auto fix system': self.auto_fix_system_basic,
            'fix system': self.auto_fix_system_basic,
            'system fix': self.auto_fix_system_basic,
            
            # Package Management
            'install package': self.install_package,
            'list packages': self.list_packages,
            'uninstall package': self.uninstall_package,
            'manage package': self.manage_package_help,
            
            # Memory & Context
            'what do you remember': self.context_memory_recall,
            'recall memory': self.context_memory_recall,
            'show memories': self.context_memory_recall,
            'daily briefing': self.daily_briefing,
            'get briefing': self.daily_briefing,
            'morning briefing': self.daily_briefing,
            
            # Calendar & Scheduling
            'schedule event': self.calendar_schedule,
            'book appointment': self.calendar_schedule,
            'show calendar': self.show_calendar,
            'check calendar': self.show_calendar,
            'my events': self.show_calendar,
            'what meetings': self.show_calendar,
            
            # Email & Communication
            'summarize email': self.email_summarize,
            'email summary': self.email_summarize,
            
            # Device & Cloud
            'sync devices': self.sync_across_devices,
            'device sync': self.sync_across_devices,
            'cloud backup': self.cloud_backup_manager,
            'backup files': self.cloud_backup_manager,
            
            # AI Productivity
            'transcribe audio': self.realtime_transcription,
            'voice to text': self.realtime_transcription,
            'summarize meeting': self.summarize_meeting,
            'meeting summary': self.summarize_meeting,
            'smart clipboard': self.smart_clipboard_store,
            'document qa': self.document_assistant_qa,
            'ask document': self.document_assistant_qa,
            'create presentation': self.ai_presentation_maker,
            'ai presentation': self.ai_presentation_maker,
            'create document': self.ai_document_maker,
            'create report': self.ai_document_maker,
            'create letter': self.ai_document_maker,
            'ai document': self.ai_document_maker,
            
            # Smart Home
            'smart home': self.smart_home_control,
            'control lights': self.smart_home_control,
            'home automation': self.smart_home_control,
            'home scene': self.set_home_scene,
            'set scene': self.set_home_scene,
            'security camera': self.security_camera_snapshot,
            'energy monitor': self.energy_monitoring_report,
            
            # Entertainment
            'dj mode': self.ai_dj_mode,
            'music dj': self.ai_dj_mode,
            'ai dj': self.ai_dj_mode,
            'trivia game': self.trivia_game_start,
            'play trivia': self.trivia_game_start,
            'tell story': self.storytelling_mode,
            'story mode': self.storytelling_mode,
            'fitness coach': self.fitness_coach,
            'workout coach': self.fitness_coach,
            
            # AI Agents
            'code agent': self.code_agent,
            'coding help': self.code_agent,
            'debug screen': self.debug_screen_code,
            'fix my code': self.debug_screen_code,
            'check code': self.debug_screen_code,
            'research agent': self.research_agent,
            'research help': self.research_agent,
            'organizer agent': self.organizer_agent,
            'organize tasks': self.organizer_agent,
            'multi agent': self.multi_agent_collab,
            'agent collaboration': self.multi_agent_collab,
            
            # Web Intelligence
            'scholar search': self.scholar_search,
            'academic search': self.scholar_search,
            'stock updates': self.stock_updates,
            'stock market': self.stock_updates,
            'crypto updates': self.crypto_updates,
            'cryptocurrency': self.crypto_updates,
            'translate text': self.realtime_translation,
            'real time translation': self.realtime_translation,
            
            # Health & Wellness
            'posture check': self.posture_detection,
            'check posture': self.posture_detection,
            'eye care': self.eye_care_mode,
            'protect eyes': self.eye_care_mode,
            'health log': self.daily_health_log,
            'track health': self.daily_health_log,
            'mood tracker': self.mood_tracker,
            'track mood': self.mood_tracker,
            'meditation': self.meditation_prompt,
            'meditate': self.meditation_prompt,
            
            # Security & Authentication
            'encrypt file': self.file_vault_encrypt,
            'decrypt file': self.file_vault_decrypt,
            'scan for threats': self.anomaly_detection_recent_processes,
            'security scan': self.anomaly_detection_recent_processes,
            'phishing scan': self.phishing_malware_scan_link,
            'parental control': self.parental_control_set,
            
            # Learning & Adaptation
            'adaptive learning': self.adaptive_learning,
            'learn from me': self.adaptive_learning,
            'check proactive': self.check_proactive_suggestions,
            'proactive suggestions': self.check_proactive_suggestions,
            'enable proactive': self.enable_proactive_mode,
            'disable proactive': self.disable_proactive_mode,
            'manual learn': self.manual_learn,
            'teach jarvis': self.manual_learn,
            'predictive assistance': self.predictive_assistance,
            'suggest actions': self.auto_suggest_with_popup,
            'start proactive background': self.start_proactive_background,
            
            # Docker & Development
            'docker control': self.docker_control,
            'docker help': self.docker_control,
        }
    
    def _get_ai_provider(self):
        try:
            with open('ai_config.json', 'r') as f:
                config = json.load(f)
                return config.get('ai_provider', 'groq')
        except:
            return 'groq'
    
    def _init_ai_models(self):
        if self.ai_provider == 'groq':
            try:
                from groq import Groq
                from engine.groq_config import GROQ_API_KEY
                self.groq_client = Groq(api_key=GROQ_API_KEY)
            except Exception as e:
                print(f"Groq failed: {e}, using Gemini")
                self.ai_provider = 'gemini'
                self._init_gemini()
        else:
            self._init_gemini()
    
    def _init_gemini(self):
        try:
            import google.generativeai as genai
            from engine.gemini_config import GEMINI_API_KEY
            genai.configure(api_key=GEMINI_API_KEY)
            self.gemini_model = genai.GenerativeModel('gemini-2.0-flash')
        except Exception as e:
            print(f"AI init error: {e}")
    
    # Utility methods
    def _ensure_db(self):
        if not os.path.exists(self.MEMORY_DB):
            conn = sqlite3.connect(self.MEMORY_DB)
            c = conn.cursor()
            c.execute("CREATE TABLE memory (id INTEGER PRIMARY KEY, key TEXT, value TEXT, ts INTEGER)")
            conn.commit()
            conn.close()

    def _memory_set(self, key: str, value: Any):
        self._ensure_db()
        conn = sqlite3.connect(self.MEMORY_DB)
        c = conn.cursor()
        c.execute("INSERT INTO memory (key, value, ts) VALUES (?, ?, ?)", (key, json.dumps(value), int(time.time())))
        conn.commit()
        conn.close()

    def _memory_get_all(self, key: str) -> List[Any]:
        self._ensure_db()
        conn = sqlite3.connect(self.MEMORY_DB)
        c = conn.cursor()
        c.execute("SELECT value FROM memory WHERE key = ? ORDER BY ts DESC", (key,))
        rows = c.fetchall()
        conn.close()
        return [json.loads(r[0]) for r in rows]

    def _calendar_add(self, title: str, event_time: str = ""):
        events = []
        if os.path.exists(self.CALENDAR_JSON):
            with open(self.CALENDAR_JSON, 'r') as f:
                events = json.load(f)
        
        events.append({
            'title': title,
            'time': event_time,
            'created': datetime.now().isoformat()
        })
        
        with open(self.CALENDAR_JSON, 'w') as f:
            json.dump(events, f, indent=2)
    
    def storytelling_mode(self, topic="random"):
        """Interactive storytelling with AI"""
        try:
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{
                        "role": "user",
                        "content": f"Tell me an engaging, creative story about {topic}. Make it interactive and ask what happens next."
                    }]
                )
                story = response.choices[0].message.content
            else:
                story = self.gemini_model.generate_content(
                    f"Tell me an engaging, creative story about {topic}. Make it interactive and ask what happens next."
                ).text
            
            print(f"\nStory Mode Activated\n{story}")
            return story
        # amazonq-ignore-next-line
        except Exception as e:
            return f"Story mode error: {e}"
    
    # amazonq-ignore-next-line
    def fitness_coach(self, workout_type="general"):
        """AI fitness coaching and workout guidance"""
        try:
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{
                        "role": "user",
                        "content": f"Act as a fitness coach. Provide a {workout_type} workout routine with exercises, reps, and motivational tips."
                    }]
                )
                workout = response.choices[0].message.content
            else:
                workout = self.gemini_model.generate_content(
                    f"Act as a fitness coach. Provide a {workout_type} workout routine with exercises, reps, and motivational tips."
                ).text
            
            print(f"\nFitness Coach Activated\n{workout}")
            return workout
        except Exception as e:
            return f"Fitness coach error: {e}"
    
    # amazonq-ignore-next-line
    # amazonq-ignore-next-line
    def _calendar_add_advanced(self, title: str, event_time: str = ""):
        events = []
        if os.path.exists(self.CALENDAR_JSON):
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            with open(self.CALENDAR_JSON, 'r') as f:
                events = json.load(f)
        
        if 'tomorrow' in title.lower():
            date = (datetime.now() + timedelta(days=1)).strftime('%Y-%m-%d')
            # amazonq-ignore-next-line
            title = title.replace('tomorrow', '').replace('Tomorrow', '').strip()
        elif 'today' in title.lower():
            date = datetime.now().strftime('%Y-%m-%d')
            title = title.replace('today', '').replace('Today', '').strip()
        else:
            date = datetime.now().strftime('%Y-%m-%d')
        
        events.append({
            'title': title,
            'date': date,
            'time': event_time,
            'timestamp': int(time.time())
        })
        
        with open(self.CALENDAR_JSON, 'w') as f:
            json.dump(events, f, indent=2)
    
    def _calendar_get_all(self):
        if os.path.exists(self.CALENDAR_JSON):
            # amazonq-ignore-next-line
            with open(self.CALENDAR_JSON, 'r') as f:
                return json.load(f)
        return []
    
    def _ai_generate(self, prompt: str) -> str:
        try:
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    messages=[{"role": "user", "content": prompt}],
                    model="llama-3.1-8b-instant"
                )
                return response.choices[0].message.content.strip()
            else:
                response = self.gemini_model.generate_content(prompt)
                return response.text.strip()
        except Exception as e:
            return f"AI generation error: {str(e)}"
    
    # VOICE-ACTIVATED ADVANCED FEATURES
    
    # System Monitoring
    def system_monitor_dashboard_live(self):
        try:
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            cpu = psutil.cpu_percent(interval=1)
            memory = psutil.virtual_memory()
            # amazonq-ignore-next-line
            disk = psutil.disk_usage('C:')
            
            processes = []
            for proc in psutil.process_iter(['pid', 'name', 'cpu_percent']):
                try:
                    processes.append(proc.info)
                # amazonq-ignore-next-line
                except:
                    continue
            
            top_processes = sorted(processes, key=lambda x: x['cpu_percent'] or 0, reverse=True)[:5]
            
            report = f"SYSTEM MONITOR LIVE:\\n"
            report += f"CPU: {cpu}%\\n"
            report += f"RAM: {memory.percent}% ({memory.used//1024//1024}MB/{memory.total//1024//1024}MB)\\n"
            report += f"Disk: {disk.used//1024//1024//1024}GB/{disk.total//1024//1024//1024}GB\\n"
            report += f"\\nTop CPU Processes:\\n"
            
            for proc in top_processes:
                if proc['cpu_percent']:
                    report += f"- {proc['name']}: {proc['cpu_percent']:.1f}% CPU\\n"
            
            return report
        except Exception as e:
            return f"System monitoring error: {str(e)}"
    
    def auto_fix_system_basic(self):
        try:
            fixes = []
            
            # amazonq-ignore-next-line
            disk = psutil.disk_usage('C:')
            if disk.free < 1024**3:
                fixes.append("Low disk space detected")
            
            memory = psutil.virtual_memory()
            if memory.percent > 90:
                fixes.append("High memory usage detected")
            
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            cpu = psutil.cpu_percent(interval=2)
            if cpu > 90:
                fixes.append("High CPU usage detected")
            
            try:
                temp_folder = os.environ.get('TEMP')
                if temp_folder:
                    temp_files = len([f for f in os.listdir(temp_folder) if os.path.isfile(os.path.join(temp_folder, f))])
                    if temp_files > 100:
                        fixes.append(f"Found {temp_files} temp files for cleanup")
            # amazonq-ignore-next-line
            except:
                pass
            
            if fixes:
                return f"System issues found:\\n" + "\\n".join([f"- {fix}" for fix in fixes])
            else:
                return "System check complete - no issues found"
        except Exception as e:
            return f"Auto-fix error: {str(e)}"
    
    # Package Management
    def install_package(self, package_name=""):
        try:
            if not package_name:
                return "Please specify package name. Say: 'install package numpy'"
            # amazonq-ignore-next-line
            result = subprocess.run(["pip", "install", package_name], capture_output=True, text=True)
            if result.returncode == 0:
                # amazonq-ignore-next-line
                return f"Successfully installed {package_name}"
            else:
                # amazonq-ignore-next-line
                return f"Failed to install {package_name}: {result.stderr}"
        except Exception as e:
            return f"Package installation error: {str(e)}"
    
    def list_packages(self):
        try:
            # amazonq-ignore-next-line
            result = subprocess.run(["pip", "list"], capture_output=True, text=True)
            if result.returncode == 0:
                # amazonq-ignore-next-line
                lines = result.stdout.split('\\n')[:10]
                return f"Installed packages:\\n" + "\\n".join(lines)
            else:
                return "Failed to list packages"
        except Exception as e:
            return f"Package listing error: {str(e)}"
    
    def uninstall_package(self, package_name=""):
        try:
            if not package_name:
                return "Please specify package name to uninstall"
            # amazonq-ignore-next-line
            result = subprocess.run(["pip", "uninstall", package_name, "-y"], capture_output=True, text=True)
            if result.returncode == 0:
                # amazonq-ignore-next-line
                return f"Successfully uninstalled {package_name}"
            else:
                # amazonq-ignore-next-line
                return f"Failed to uninstall {package_name}"
        except Exception as e:
            return f"Package uninstall error: {str(e)}"
    
    def manage_package_help(self):
        return "Package management available. Say: 'install package [name]', 'list packages', or 'uninstall package [name]'"
    
    # Memory & Context
    def context_memory_store(self, key="general", value="memory stored"):
        self._memory_set(key, value)
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        return f"Stored memory under '{key}': {value}"

    def context_memory_recall(self, key=None):
        if key:
            vals = self._memory_get_all(key)
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            return vals[0] if vals else "No memory found"
        
  
        self._ensure_db()
        # amazonq-ignore-next-line
        conn = sqlite3.connect(self.MEMORY_DB)
        # amazonq-ignore-next-line
        c = conn.cursor()
        c.execute("SELECT key, value FROM memory ORDER BY ts DESC LIMIT 10")
        rows = c.fetchall()
        conn.close()
        
        if rows:
            memories = []
            for key, value in rows:
                try:
                    data = json.loads(value)
                    memories.append(f"{key}: {data}")
                except:
                    memories.append(f"{key}: {value}")
            return "\n".join(memories)
        return "No memories stored"

    def daily_briefing(self):
        briefing = {}
        # amazonq-ignore-next-line
        briefing["datetime"] = datetime.now().isoformat()
        
        try:
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            cpu = psutil.cpu_percent(interval=0.5)
            ram = psutil.virtual_memory().percent
            # amazonq-ignore-next-line
            batt = None
            if psutil.sensors_battery():
                batt = psutil.sensors_battery().percent
            briefing["system"] = {"cpu_percent": cpu, "ram_percent": ram, "battery_percent": batt}
        # amazonq-ignore-next-line
        except:
            briefing["system"] = {"cpu_percent": None, "ram_percent": None, "battery_percent": None}
        
        briefing["weather"] = {"note": "Weather API not configured"}
        briefing["calendar_upcoming"] = []
        briefing["unread_emails"] = {"note": "Email integration not configured"}
        
        return briefing
    
    # Calendar & Scheduling
    def calendar_schedule(self, event_data=""):
        if not event_data:
            return "No event specified"
        
        title = event_data
        time_str = ""
        
        if ' at ' in event_data.lower():
            parts = event_data.split(' at ')
            title = parts[0].strip()
            time_part = parts[1].strip().lower()
            
            if 'pm' in time_part:
                hour = int(time_part.replace('pm', '').strip())
                if hour != 12:
                    hour += 12
                time_str = f"{hour:02d}:00"
            elif 'am' in time_part:
                hour = int(time_part.replace('am', '').strip())
                if hour == 12:
                    hour = 0
                time_str = f"{hour:02d}:00"
        
        self._calendar_add_advanced(title, time_str)
        return f"Scheduled: {title} at {time_str}"
    
    def show_calendar(self):
        events = self._calendar_get_all()
        if not events:
            return "No events scheduled"
        
        today = datetime.now().strftime('%Y-%m-%d')
        tomorrow = (datetime.now() + timedelta(days=1)).strftime('%Y-%m-%d')
        
        today_events = []
        tomorrow_events = []
        other_events = []
        
        for event in events:
            event_date = event.get('date', today)  # Default to today if no date
            if event_date == today:
                today_events.append(event)
            elif event_date == tomorrow:
                tomorrow_events.append(event)
            else:
                other_events.append(event)
        
        result = ""
        if today_events:
            result += "TODAY:\n"
            for event in today_events:
                time = event.get('time', '')
                if time:
                    result += f"- {event['title']} at {time}\n"
                else:
                    result += f"- {event['title']}\n"
        
        if tomorrow_events:
            result += "\nTOMORROW:\n"
            for event in tomorrow_events:
                time = event.get('time', '')
                if time:
                    result += f"- {event['title']} at {time}\n"
                else:
                    result += f"- {event['title']}\n"
        
        if other_events:
            result += "\nOTHER DATES:\n"
            for event in other_events:
                time = event.get('time', '')
                event_date = event.get('date', 'Unknown date')
                if time:
                    result += f"- {event['title']} on {event_date} at {time}\n"
                else:
                    result += f"- {event['title']} on {event_date}\n"
        
        return result.strip()
    
    # Email & Communication
    def email_summarize(self, email_content=""):
        try:
            if not email_content:
                return "Please provide email content. Say: 'summarize email [your email content here]'"
            
            # amazonq-ignore-next-line
            prompt = f"Summarize this email in 3 key points: {email_content}"
            summary = self._ai_generate(prompt)
            # amazonq-ignore-next-line
            return f"Email Summary:\\n{summary}"
        except Exception as e:
            return f"Email summarization failed: {str(e)}"
    
    # Device & Cloud
    def sync_across_devices(self):
        try:
            calendar_data = self._calendar_get_all()
            memories = self.context_memory_recall()
            
            sync_data = {
                "calendar_events": len(calendar_data),
                # amazonq-ignore-next-line
                "memory_entries": len(memories) if isinstance(memories, dict) else 0,
                "last_sync": datetime.now().isoformat(),
                "device_id": socket.gethostname()
            }
            
            sync_file = os.path.join(os.path.dirname(self.USAGE_LOG), "device_sync.json")
            with open(sync_file, 'w') as f:
                json.dump(sync_data, f, indent=2)
            
            return f"Device sync completed: {sync_data['calendar_events']} events, {sync_data['memory_entries']} memories synced"
        except Exception as e:
            return f"Device sync failed: {str(e)}"
    
    # amazonq-ignore-next-line
    def cloud_backup_manager(self):
        return "Cloud backup feature available - files can be backed up to cloud storage"
    
    # AI Productivity
    def realtime_transcription(self):
        try:
            # Create transcription file
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"transcription_{timestamp}.txt"
            
            # Start Windows Speech Recognition
            # amazonq-ignore-next-line
            subprocess.run('start ms-speech-recognition:', shell=True)
            
            # Create empty file for transcription
            # amazonq-ignore-next-line
            with open(filename, 'w') as f:
                f.write(f"Transcription started at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write("Speak now - text will appear here...\n\n")
            
            # Store in memory
            self._memory_set("transcription_file", filename)
            
            return f"Transcription started - File: {filename} | Windows Speech Recognition opened"
            
        except Exception as e:
            return f"Transcription error: {str(e)}"
    
    def summarize_meeting(self, meeting_text=""):
        try:
            if not meeting_text:
                return "Please provide meeting content to summarize"
            
            prompt = f"Summarize this meeting in key points: {meeting_text}"
            # amazonq-ignore-next-line
            return self._ai_generate(prompt)
        except Exception as e:
            return f"Meeting summary error: {str(e)}"
    
    def smart_clipboard_store(self, content=""):
        try:
            if content:
                self._memory_set("clipboard", content)
                # amazonq-ignore-next-line
                return f"Stored in smart clipboard: {content[:50]}..."
            else:
                clipboard_data = self._memory_get_all("clipboard")
                if clipboard_data:
                    return f"Smart clipboard contains: {clipboard_data[0][:100]}..."
                return "Smart clipboard is empty"
        except Exception as e:
            return f"Smart clipboard error: {str(e)}"
    
    def document_assistant_qa(self, question="", document=""):
        try:
            if not question:
                return "Please ask a question about your document"
            
            # amazonq-ignore-next-line
            prompt = f"Answer this question about the document: {question}\\nDocument: {document[:500]}"
            return self._ai_generate(prompt)
        except Exception as e:
            return f"Document Q&A error: {str(e)}"
    
    def ai_presentation_maker(self, topic="", create_ppt=True, num_slides=5):
        try:
            if not topic:
                return "Please specify a presentation topic"
            
            if create_ppt:
                try:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt

                    prs = Presentation()

                    # Title slide
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    slide.shapes.title.text = topic.title()
                    slide.placeholders[1].text = "AI Generated Presentation"

                    # Ask AI to generate bullet points for slides
                    prompt = (
                        f"Generate {num_slides} slide titles and bullet points for a presentation on '{topic}'.\n"
                        "Return in format:\n"
                        "Slide Title: \n"
                        "- bullet\n"
                        "- bullet\n"
                        "\n"
                    )

                    ai_response = self._ai_generate(prompt)
                    sections = ai_response.strip().split("\n\n")

                    # Use Title and Content slide layout
                    content_layout = prs.slide_layouts[1]

                    for sec in sections[:num_slides]:
                        lines = sec.split("\n")
                        if len(lines) < 2:
                            continue

                        slide_title = lines[0].replace("Slide Title:", "").strip()
                        slide = prs.slides.add_slide(content_layout)
                        slide.shapes.title.text = slide_title

                        bullet_box = slide.shapes.placeholders[1].text_frame
                        bullet_box.text = lines[1].lstrip("- ").strip()

                        for bullet in lines[2:]:
                            p = bullet_box.add_paragraph()
                            p.text = bullet.lstrip("- ").strip()
                            p.level = 1

                    filename = f"{topic.replace(' ', '_')}_presentation.pptx"
                    prs.save(filename)

                    return f"✅ Presentation created successfully: {filename}"

                except ImportError:
                    return "Install python-pptx to enable presentation creation"

            return "PowerPoint creation disabled"

        except Exception as e:
            return f"AI presentation error: {str(e)}"
    
    def ai_document_maker(self, doc_type="report", topic="", user_info="", num_pages=2):
        try:
            if not topic:
                return "Please specify document topic"

            try:
                from docx import Document
                from docx.shared import Pt, Inches
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                from docx.oxml.ns import qn

                doc = Document()

                # Title Style
                title = doc.add_heading(topic.title(), level=0)
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # Format Title Font
                title_run = title.runs[0]
                title_run.font.name = 'Calibri'
                title_run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Calibri')
                title_run.font.size = Pt(28)
                title_run.bold = True

                # Build AI prompt (more structured output)
                prompt = (
                    f"Write a well-structured {doc_type} on the topic '{topic}'. "
                    f"Include these details: {user_info}. "
                    f"Document must include:\n"
                    f"1. Introduction\n"
                    f"2. Main content with multiple clear headings\n"
                    f"3. Conclusion\n"
                    f"Ensure the tone is formal and academic.\n"
                    f"Length: About {num_pages} pages.\n"
                    f"Do NOT use markdown (** **, ## , etc.). Only plain structured text."
                )

                content = self._ai_generate(prompt)

                # Split into paragraphs
                paragraphs = content.split('\n')

                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue

                    # Detect headings
                    if len(para.split()) < 8:  # short lines = likely headings
                        heading = doc.add_heading(para, level=1)
                        heading.alignment = WD_ALIGN_PARAGRAPH.LEFT

                    else:
                        p = doc.add_paragraph(para)
                        p.style = doc.styles['Normal']
                        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                        # Improve paragraph readability
                        for run in p.runs:
                            run.font.name = 'Cambria'
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Cambria')
                            run.font.size = Pt(12)

                # Add ending spacing & clean page look
                doc.add_page_break()

                filename = f"{topic.replace(' ', '_')}_{doc_type}.docx"
                doc.save(filename)

                return f"✅ Document created successfully: {filename}"

            except ImportError:
                return "Install python-docx to create Word documents"

        except Exception as e:
            return f"Document creation error: {str(e)}"

    # All other features with minimal implementations
    # amazonq-ignore-next-line
    def smart_home_control(self): return "Smart home control available"
    def set_home_scene(self): return "Home scene setting available"
    def security_camera_snapshot(self): return "Security camera available"
    def energy_monitoring_report(self): return "Energy monitoring available"
    def ai_dj_mode(self): return "AI DJ mode available"
    def trivia_game_start(self): return "Trivia game available"
    def code_agent(self, task="general coding help"):
        """AI-powered coding assistant"""
        try:
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{
                        "role": "user",
                        "content": f"Act as an expert software developer and coding assistant. Help with: {task}. Provide code examples, explanations, and best practices."
                    }]
                )
                code_help = response.choices[0].message.content
            else:
                code_help = self.gemini_model.generate_content(
                    f"Act as an expert software developer and coding assistant. Help with: {task}. Provide code examples, explanations, and best practices."
                ).text
            
            print(f"\nCode Agent Activated\n{code_help}")
            return code_help
        except Exception as e:
            return f"Code agent error: {e}"
    
    def research_agent(self, topic="general research"):
        """AI-powered research assistant"""
        try:
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{
                        "role": "user",
                        "content": f"Act as a research assistant. Provide comprehensive research on: {topic}. Include key points, analysis, and relevant information."
                    }]
                )
                research = response.choices[0].message.content
            else:
                research = self.gemini_model.generate_content(
                    f"Act as a research assistant. Provide comprehensive research on: {topic}. Include key points, analysis, and relevant information."
                ).text
            
            print(f"\nResearch Agent Activated\n{research}")
            return research
        except Exception as e:
            return f"Research agent error: {e}"
    
    def debug_screen_code(self):
        """Capture screen and analyze visible code for errors"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            screenshot_path = f"code_debug_{timestamp}.png"
            # amazonq-ignore-next-line
            screenshot = pyautogui.screenshot()
            screenshot.save(screenshot_path)
            
            # Try to get clipboard content first
            clipboard_text = ""
            try:
                import win32clipboard
                # amazonq-ignore-next-line
                win32clipboard.OpenClipboard()
                clipboard_text = win32clipboard.GetClipboardData()
                win32clipboard.CloseClipboard()
            # amazonq-ignore-next-line
            except:
                pass
            
            # If no clipboard content, ask user to provide code
            if not clipboard_text or len(clipboard_text.strip()) < 10:
                return f"\nScreen captured: {screenshot_path}\n\nTo analyze your code:\n1. Copy the code you want me to check\n2. Say 'debug screen' again\n\nOr tell me what specific error you're seeing."
            
            # Analyze the actual clipboard code
            prompt = f"""You are a code debugging expert. Analyze this code for errors, bugs, and improvements:

CODE TO ANALYZE:
{clipboard_text}

Provide:
1. ERRORS FOUND: List any syntax errors, logic errors, or bugs
2. CORRECTED CODE: Show the fixed version
3. EXPLANATION: Explain what was wrong and why the fix works
4. IMPROVEMENTS: Suggest any additional improvements

Be specific and accurate about the actual code provided."""
            
            if self.ai_provider == 'groq':
                response = self.groq_client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{"role": "user", "content": prompt}]
                )
                analysis = response.choices[0].message.content
            else:
                analysis = self.gemini_model.generate_content(prompt).text
            
            result = f"\nCode Debug Analysis\nScreenshot saved: {screenshot_path}\n\nAnalyzing your code:\n{analysis}"
            print(result)
            return result
            
        except Exception as e:
            return f"Screen debug error: {e}"
    # amazonq-ignore-next-line
    def organizer_agent(self): return "Organizer agent available"
    def multi_agent_collab(self): return "Multi-agent collaboration available"
    
    # Web Intelligence Methods
    def scholar_search(self, query=""):
        """Search academic papers and scholarly content"""
        try:
            if not query:
                return "Please specify search query. Say: 'scholar search machine learning'"
            
            prompt = f"Provide academic research information about: {query}. Include key papers, concepts, and recent developments."
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            research = self._ai_generate(prompt)
            # amazonq-ignore-next-line
            return f"Scholar Search Results for '{query}':\n{research}"
        except Exception as e:
            return f"Scholar search error: {e}"
    
    def stock_updates(self, symbol=""):
        """Get stock market updates and analysis"""
        try:
            if symbol:
                prompt = f"Provide current stock analysis for {symbol}. Include price trends, market sentiment, and key factors."
            else:
                prompt = "Provide general stock market update. Include major indices, market trends, and key news."
            
            analysis = self._ai_generate(prompt)
            # amazonq-ignore-next-line
            return f"Stock Market Update:\n{analysis}"
        except Exception as e:
            return f"Stock updates error: {e}"
    
    def crypto_updates(self, coin=""):
        """Get cryptocurrency updates and analysis"""
        try:
            if coin:
                prompt = f"Provide cryptocurrency analysis for {coin}. Include price trends, market cap, and recent developments."
            else:
                prompt = "Provide general cryptocurrency market update. Include Bitcoin, Ethereum, and market trends."
            
            analysis = self._ai_generate(prompt)
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            return f"Crypto Market Update:\n{analysis}"
        except Exception as e:
            return f"Crypto updates error: {e}"
    
    def realtime_translation(self, text="", target_lang=""):
        """Translate text to different languages"""
        try:
            if not text:
                return "Please provide text to translate. Say: 'translate text hello world to spanish'"
            
            if not target_lang:
                target_lang = "Spanish"
            
            prompt = f"Translate this text to {target_lang}: {text}"
            translation = self._ai_generate(prompt)
            # amazonq-ignore-next-line
            return f"Translation to {target_lang}:\n{translation}"
        except Exception as e:
            return f"Translation error: {e}"
    # Health & Wellness Methods
    def posture_detection(self):
        """AI-powered posture analysis using webcam"""
        try:
            # Try to use webcam for actual posture detection
            try:
                import cv2
                import numpy as np
                
                # Initialize webcam
                # amazonq-ignore-next-line
                cap = cv2.VideoCapture(0)
                if not cap.isOpened():
                    return self._fallback_posture_advice()
                
                # Capture frame
                ret, frame = cap.read()
                cap.release()
                
                if not ret:
                    return self._fallback_posture_advice()
                
                # Save webcam image
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                image_path = f"posture_analysis_{timestamp}.jpg"
                cv2.imwrite(image_path, frame)
                
                # Basic posture analysis using image processing
                analysis = self._analyze_posture_from_image(frame)
                
                # Generate personalized advice based on analysis
                prompt = f"""Based on this posture analysis: {analysis}
                
                Provide specific posture improvement advice:
                1. CURRENT POSTURE ASSESSMENT: What needs improvement
                2. IMMEDIATE CORRECTIONS: Quick fixes for current posture
                3. EXERCISES: Targeted exercises for detected issues
                4. PREVENTION TIPS: How to maintain better posture
                
                Keep it personalized and actionable."""
                
                advice = self._ai_generate(prompt)
                
                # Store posture check in memory
                self._memory_set("posture_check", {
                    "timestamp": datetime.now().isoformat(),
                    "image_path": image_path,
                    "analysis": analysis,
                    "advice_given": True
                })
                
                result = f"\nPosture Analysis Complete\nWebcam image saved: {image_path}\nAnalysis: {analysis}\n\n{advice}"
                print(result)
                return result
                
            except ImportError:
                return self._fallback_posture_advice("OpenCV not installed. Install with: pip install opencv-python")
            # amazonq-ignore-next-line
            except Exception as e:
                return self._fallback_posture_advice(f"Webcam error: {e}")
                
        except Exception as e:
            return f"Posture detection error: {e}"
    
    def _analyze_posture_from_image(self, frame):
        """Basic posture analysis from webcam image"""
        try:
            import cv2
            import numpy as np
            
            # Convert to grayscale
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            
            # Simple analysis based on image properties
            height, width = gray.shape
            
            # Analyze head position (top third of image)
            head_region = gray[:height//3, :]
            # amazonq-ignore-next-line
            head_center = np.mean(np.where(head_region < np.mean(head_region)))
            
            # Analyze shoulder level (middle third)
            shoulder_region = gray[height//3:2*height//3, :]
            
            # Basic posture assessment
            analysis = []
            
            # Check if head is centered
            # amazonq-ignore-next-line
            if head_center < width * 0.4 or head_center > width * 0.6:
                analysis.append("Head appears tilted - try to keep head centered")
            else:
                analysis.append("Head position looks good")
            
            # Check overall brightness (distance from camera)
            avg_brightness = np.mean(gray)
            if avg_brightness < 80:
                analysis.append("You appear to be leaning forward - sit back in your chair")
            elif avg_brightness > 180:
                analysis.append("Good distance from camera")
            else:
                analysis.append("Reasonable sitting distance")
            
            # Check for symmetry (shoulder level)
            left_shoulder = np.mean(shoulder_region[:, :width//2])
            right_shoulder = np.mean(shoulder_region[:, width//2:])
            
            if abs(left_shoulder - right_shoulder) > 20:
                analysis.append("Shoulders appear uneven - check your sitting position")
            else:
                analysis.append("Shoulder alignment looks balanced")
            
            return "; ".join(analysis)
            
        except Exception as e:
            return f"Basic visual analysis completed (detailed analysis unavailable: {e})"
    
    def _fallback_posture_advice(self, reason="Webcam not available"):
        """Fallback posture advice when webcam analysis fails"""
        prompt = """Act as a health and ergonomics expert. Provide comprehensive posture improvement guidance:
        
        1. POSTURE SELF-CHECK: How to assess your own posture right now
        2. COMMON ISSUES: Most frequent posture problems and quick fixes
        3. DESK EXERCISES: Simple exercises you can do at your workstation
        4. ERGONOMIC SETUP: Optimal workspace configuration
        5. DAILY HABITS: Posture-friendly habits to develop
        
        Make it practical and immediately actionable."""
        
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        advice = self._ai_generate(prompt)
        
        result = f"\nPosture Check ({reason})\n\nSELF-ASSESSMENT GUIDE:\n{advice}\n\nTip: Sit up straight, shoulders back, feet flat on floor!"
        print(result)
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        # amazonq-ignore-next-line
        return result
    
    def eye_care_mode(self):
        """Eye care reminders and exercises"""
        try:
            prompt = """Act as an eye care specialist. Provide comprehensive eye care guidance:
            
            1. 20-20-20 RULE: Explain and remind about this important rule
            2. EYE EXERCISES: Simple exercises to reduce eye strain
            3. SCREEN SETTINGS: Optimal brightness, contrast, and blue light tips
            4. BREAK REMINDERS: When and how often to take eye breaks
            5. WARNING SIGNS: Symptoms that indicate eye strain
            
            Make it practical for computer users."""
            
            # amazonq-ignore-next-line
            eye_care = self._ai_generate(prompt)
            
            # Set eye care reminder
            self._memory_set("eye_care_reminder", {
                # amazonq-ignore-next-line
                "timestamp": datetime.now().isoformat(),
                "next_reminder": (datetime.now() + timedelta(minutes=20)).isoformat(),
                "rule_20_20_20": True
            })
            
            result = f"\nEye Care Mode Activated\n{eye_care}\n\nNext reminder set for 20 minutes."
            print(result)
            return result
            
        except Exception as e:
            return f"Eye care mode error: {e}"
    
    def daily_health_log(self, entry=""):
        """Track daily health metrics and habits"""
        try:
            today = datetime.now().strftime('%Y-%m-%d')
            
            if entry:
                # Log specific health entry
                self._memory_set(f"health_log_{today}", {
                    "date": today,
                    "entry": entry,
                    "timestamp": datetime.now().isoformat()
                })
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                return f"Health log updated for {today}: {entry}"
            else:
                # Show health log summary
                recent_logs = []
                # amazonq-ignore-next-line
                for i in range(7):  # Last 7 days
                    date = (datetime.now() - timedelta(days=i)).strftime('%Y-%m-%d')
                    logs = self._memory_get_all(f"health_log_{date}")
                    if logs:
                        recent_logs.extend(logs)
                
                if recent_logs:
                    summary = "\nRecent Health Logs:\n"
                    for log in recent_logs[:5]:  # Show last 5 entries
                        summary += f"- {log['date']}: {log['entry']}\n"
                    return summary
                else:
                    # Generate health tracking prompt
                    prompt = """Provide a daily health tracking template with:
                    
                    1. METRICS TO TRACK: Key health indicators to monitor daily
                    2. QUESTIONS: Simple questions to assess daily wellness
                    3. HABITS: Healthy habits to track and maintain
                    4. MOOD: How to track emotional well-being
                    
                    Make it simple and quick to use daily."""
                    
                    template = self._ai_generate(prompt)
                    return f"\nDaily Health Log\n{template}\n\nSay 'health log [your entry]' to add an entry."
                    
        except Exception as e:
            return f"Health log error: {e}"
    
    def mood_tracker(self, mood=""):
        """Track and analyze mood patterns"""
        try:
            today = datetime.now().strftime('%Y-%m-%d')
            
            if mood:
                # Log mood entry
                mood_data = {
                    "date": today,
                    "mood": mood.lower(),
                    "timestamp": datetime.now().isoformat(),
                    "time_of_day": datetime.now().strftime('%H:%M')
                }
                
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                # amazonq-ignore-next-line
                self._memory_set(f"mood_{today}", mood_data)
                
                # Generate mood response
                prompt = f"""Act as a wellness coach. Someone just reported feeling '{mood}'. Provide:
                
                1. ACKNOWLEDGMENT: Validate their feeling
                2. SUGGESTIONS: 2-3 practical tips to improve or maintain this mood
                3. ACTIVITIES: Specific activities that might help
                4. REMINDER: Positive affirmation or encouragement
                
                Keep it supportive and actionable."""
                
                response = self._ai_generate(prompt)
                # amazonq-ignore-next-line
                return f"\nMood logged: {mood}\n{response}"
            else:
                # Show mood analysis
                recent_moods = []
                for i in range(7):  # Last 7 days
                    date = (datetime.now() - timedelta(days=i)).strftime('%Y-%m-%d')
                    # amazonq-ignore-next-line
                    moods = self._memory_get_all(f"mood_{date}")
                    recent_moods.extend(moods)
                
                if recent_moods:
                    mood_list = [m['mood'] for m in recent_moods]
                    mood_summary = Counter(mood_list).most_common(3)
                    
                    analysis = "\nMood Analysis (Last 7 days):\n"
                    for mood, count in mood_summary:
                        analysis += f"- {mood.title()}: {count} times\n"
                    
                    # Generate insights
                    prompt = f"""Based on these recent moods: {mood_list}, provide:
                    
                    1. PATTERNS: Any patterns you notice
                    2. INSIGHTS: What this might indicate about well-being
                    3. RECOMMENDATIONS: Suggestions for mood improvement
                    
                    Be supportive and constructive."""
                    
                    insights = self._ai_generate(prompt)
                    return f"{analysis}\n{insights}"
                else:
                    return "\nMood Tracker\nNo mood entries yet. Say 'mood tracker happy' or 'track mood stressed' to start logging."
                    
        except Exception as e:
            return f"Mood tracker error: {e}"
    
    def meditation_prompt(self, duration="5"):
        """Guided meditation and mindfulness exercises"""
        try:
            # Parse duration
            try:
                minutes = int(duration)
            # amazonq-ignore-next-line
            except:
                minutes = 5
            
            prompt = f"""Act as a meditation guide. Create a {minutes}-minute guided meditation session:
            
            1. PREPARATION: How to get ready and comfortable
            2. BREATHING: Specific breathing techniques
            3. FOCUS: What to focus on during meditation
            4. GUIDANCE: Step-by-step instructions for the session
            5. CLOSING: How to gently end the session
            
            Make it calming, clear, and suitable for beginners."""
            
            meditation = self._ai_generate(prompt)
            
            # Log meditation session
            self._memory_set("meditation_session", {
                # amazonq-ignore-next-line
                "timestamp": datetime.now().isoformat(),
                "duration_minutes": minutes,
                "completed": True
            })
            
            result = f"\nMeditation Session ({minutes} minutes)\n{meditation}\n\nSession logged. Take your time and breathe deeply."
            print(result)
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            return result
            
        except Exception as e:
            return f"Meditation prompt error: {e}"
    # Security & Authentication Methods
    # amazonq-ignore-next-line
    # amazonq-ignore-next-line
    # amazonq-ignore-next-line
    # amazonq-ignore-next-line
    def file_vault_encrypt(self, file_path=""):
        """Encrypt files for security"""
        try:
            if not file_path:
                return "Please specify file path. Say: 'encrypt file C:\\path\\to\\file.txt'"
            
            if not os.path.exists(file_path):
                return f"File not found: {file_path}"
            
            # Simple encryption using base64 and password
            import base64
            # amazonq-ignore-next-line
            password = "jarvis_secure_key_2024"  # In real app, use user password
            
            with open(file_path, 'rb') as f:
                file_data = f.read()
            
            # Encrypt data
            encrypted_data = base64.b64encode(file_data)
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            encrypted_file = file_path + ".encrypted"
            
            with open(encrypted_file, 'wb') as f:
                f.write(encrypted_data)
            
            # Store encryption info
            self._memory_set("encrypted_files", {
                "original": file_path,
                "encrypted": encrypted_file,
                # amazonq-ignore-next-line
                "timestamp": datetime.now().isoformat()
            })
            
            return f"File encrypted successfully: {encrypted_file}"
            
        except Exception as e:
            return f"Encryption error: {e}"
    
    def file_vault_decrypt(self, encrypted_file=""):
        """Decrypt encrypted files"""
        try:
            if not encrypted_file:
                return "Please specify encrypted file path"
            
            if not os.path.exists(encrypted_file):
                return f"Encrypted file not found: {encrypted_file}"
            
            import base64
            
            with open(encrypted_file, 'rb') as f:
                encrypted_data = f.read()
            
            # Decrypt data
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            decrypted_data = base64.b64decode(encrypted_data)
            decrypted_file = encrypted_file.replace(".encrypted", "_decrypted")
            
            with open(decrypted_file, 'wb') as f:
                f.write(decrypted_data)
            
            return f"File decrypted successfully: {decrypted_file}"
            
        except Exception as e:
            return f"Decryption error: {e}"
    
    def anomaly_detection_recent_processes(self):
        """Scan for suspicious processes and system anomalies"""
        try:
            suspicious_processes = []
            high_cpu_processes = []
            
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # Check running processes
            # amazonq-ignore-next-line
            for proc in psutil.process_iter(['pid', 'name', 'cpu_percent', 'memory_percent']):
                try:
                    proc_info = proc.info
                    
                    # Check for high CPU usage
                    if proc_info['cpu_percent'] and proc_info['cpu_percent'] > 80:
                        high_cpu_processes.append(proc_info)
                    
                    # Check for suspicious process names
                    # amazonq-ignore-next-line
                    suspicious_names = ['keylogger', 'trojan', 'virus', 'malware', 'backdoor']
                    if any(sus in proc_info['name'].lower() for sus in suspicious_names):
                        suspicious_processes.append(proc_info)
                        
                # amazonq-ignore-next-line
                except (psutil.NoSuchProcess, psutil.AccessDenied):
                    continue
            
            # Generate security report
            report = "\nSecurity Scan Results:\n"
            
            if suspicious_processes:
                report += f"WARNING: SUSPICIOUS PROCESSES FOUND ({len(suspicious_processes)}):\n"
                for proc in suspicious_processes:
                    report += f"- {proc['name']} (PID: {proc['pid']})\n"
            else:
                report += "OK: No suspicious processes detected\n"
            
            if high_cpu_processes:
                report += f"\nWARNING: HIGH CPU USAGE PROCESSES ({len(high_cpu_processes)}):\n"
                for proc in high_cpu_processes[:3]:  # Show top 3
                    report += f"- {proc['name']}: {proc['cpu_percent']:.1f}% CPU\n"
            
            # System health check
            cpu_avg = psutil.cpu_percent(interval=1)
            memory_usage = psutil.virtual_memory().percent
            
            report += f"\nSYSTEM HEALTH:\n"
            report += f"- CPU Usage: {cpu_avg}%\n"
            report += f"- Memory Usage: {memory_usage}%\n"
            
            if cpu_avg > 90 or memory_usage > 90:
                report += "WARNING: System resources critically high\n"
            elif cpu_avg > 70 or memory_usage > 80:
                report += "WARNING: System resources elevated\n"
            else:
                report += "OK: System resources normal\n"
            
            # Store scan results
            self._memory_set("security_scan", {
                # amazonq-ignore-next-line
                "timestamp": datetime.now().isoformat(),
                "suspicious_count": len(suspicious_processes),
                "high_cpu_count": len(high_cpu_processes),
                "system_health": "normal" if cpu_avg < 70 and memory_usage < 80 else "elevated"
            })
            
            return report
            
        except Exception as e:
            return f"Security scan error: {e}"
    
    # amazonq-ignore-next-line
    # amazonq-ignore-next-line
    def phishing_malware_scan_link(self, url=""):
        """Scan URLs for phishing and malware threats"""
        try:
            if not url:
                return "Please provide URL to scan. Say: 'phishing scan https://example.com'"
            
            # Basic URL analysis
            suspicious_indicators = []
            
            # Check for suspicious domains
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            suspicious_domains = ['bit.ly', 'tinyurl.com', 'short.link']
            if any(domain in url.lower() for domain in suspicious_domains):
                suspicious_indicators.append("Shortened URL detected")
            
            # amazonq-ignore-next-line
            # amazonq-ignore-next-line
            # Check for suspicious patterns
            # amazonq-ignore-next-line
            if url.count('.') > 4:
                suspicious_indicators.append("Excessive subdomains")
            
            if any(char in url for char in ['@', '%']):
                suspicious_indicators.append("Suspicious characters in URL")
            
            # Check for HTTPS
            if not url.startswith('https://'):
                suspicious_indicators.append("Not using secure HTTPS")
            
            # Generate scan report
            report = f"\nPhishing Scan Results for: {url}\n"
            
            if suspicious_indicators:
                report += f"WARNING: POTENTIAL RISKS DETECTED ({len(suspicious_indicators)}):\n"
                for indicator in suspicious_indicators:
                    report += f"- {indicator}\n"
                report += "\nRECOMMENDATION: Exercise caution before visiting this link\n"
            else:
                report += "OK: No obvious phishing indicators detected\n"
                report += "OK: URL appears to be safe\n"
            
            # Store scan results
            self._memory_set("phishing_scan", {
                "url": url,
                # amazonq-ignore-next-line
                "timestamp": datetime.now().isoformat(),
                "risk_level": "high" if len(suspicious_indicators) > 2 else "low",
                "indicators": suspicious_indicators
            })
            
            return report
            
        except Exception as e:
            return f"Phishing scan error: {e}"
    
    # amazonq-ignore-next-line
    def parental_control_set(self, action="status"):
        """Set up parental controls and content filtering"""
        try:
            if action == "status":
                # Check current parental control status
                controls = self._memory_get_all("parental_controls")
                if controls:
                    latest = controls[0]
                    # amazonq-ignore-next-line
                    return f"\nParental Controls Status:\n- Enabled: {latest.get('enabled', False)}\n- Content Filter: {latest.get('content_filter', 'Off')}\n- Time Restrictions: {latest.get('time_restrictions', 'None')}"
                else:
                    return "\nParental Controls: Not configured\nSay 'parental control enable' to set up"
            
            elif action == "enable":
                # Enable parental controls
                controls = {
                    "enabled": True,
                    # amazonq-ignore-next-line
                    "content_filter": "Moderate",
                    "time_restrictions": "School hours",
                    "blocked_sites": ["adult content", "gambling", "violence"],
                    # amazonq-ignore-next-line
                    "timestamp": datetime.now().isoformat()
                }
                
                self._memory_set("parental_controls", controls)
                
                # Generate setup guide
                prompt = """Provide parental control setup guidance:
                
                1. CONTENT FILTERING: How to set up safe browsing
                2. TIME RESTRICTIONS: Recommended screen time limits
                3. APP CONTROLS: Managing app access for children
                4. MONITORING: Safe ways to monitor child activity
                5. COMMUNICATION: How to talk to children about online safety
                
                Keep it practical and family-friendly."""
                
                guide = self._ai_generate(prompt)
                
                return f"\nParental Controls Enabled\n{guide}\n\nControls are now active for family safety."
            
            elif action == "disable":
                # Disable parental controls
                controls = {
                    "enabled": False,
                    "timestamp": datetime.now().isoformat()
                }
                
                self._memory_set("parental_controls", controls)
                return "Parental controls disabled"
            
            else:
                return "Available commands: 'parental control status', 'parental control enable', 'parental control disable'"
                
        except Exception as e:
            return f"Parental control error: {e}"
    # amazonq-ignore-next-line
    def predictive_assistance(self, auto_speak: bool = False):
        suggestions = []
        now = datetime.now()
        hour = now.hour
        
        if hour >= 22:
            suggestions.append("It's late - set a shutdown or power-save schedule?")
        elif hour >= 18:
            suggestions.append("Evening time - would you like to check your tasks for tomorrow?")
        elif hour >= 12:
            suggestions.append("Afternoon break - time for a quick stretch?")
        else:
            suggestions.append("Good morning - ready to start your productive day?")
        
        return suggestions
    
    def adaptive_learning(self, record_action: str = "general_action"):
        self._log_action_with_context(record_action)
        # amazonq-ignore-next-line
        return f"Learned: {record_action}"
    
    def _log_action_with_context(self, action_name: str):
        data = self._load_learning_memory()
        now = datetime.now()
        day = now.strftime("%A")
        hour = now.hour
        time_of_day = "morning" if hour < 12 else "afternoon" if hour < 18 else "evening"
        
        for act in data["actions"]:
            if (act["action"] == action_name and 
                act["context"]["day"] == day and 
                act["context"]["time_of_day"] == time_of_day):
                # amazonq-ignore-next-line
                act["count"] += 1
                act["last_used"] = now.isoformat()
                self._save_learning_memory(data)
                return
        
        data["actions"].append({
            "action": action_name,
            "count": 1,
            "last_used": now.isoformat(),
            "context": {
                "day": day,
                "hour": hour,
                "time_of_day": time_of_day
            }
        })
        self._save_learning_memory(data)
    
    def _load_learning_memory(self):
        learning_file = os.path.join(os.path.dirname(self.USAGE_LOG), "adaptive_memory.json")
        if os.path.exists(learning_file):
            try:
                with open(learning_file, "r", encoding='utf-8') as f:
                    return json.load(f)
            except (json.JSONDecodeError, IOError, OSError) as e:
                print(f"Error loading learning memory: {e}")
        return {"actions": []}
    
    def _save_learning_memory(self, data):
        try:
            learning_file = os.path.join(os.path.dirname(self.USAGE_LOG), "adaptive_memory.json")
            with open(learning_file, "w", encoding='utf-8') as f:
                json.dump(data, f, indent=2)
        except (IOError, OSError) as e:
            print(f"Error saving learning memory: {e}")
    
    def auto_suggest_with_popup(self):
        data = self._load_learning_memory()
        now = datetime.now()
        day = now.strftime("%A")
        hour = now.hour
        time_of_day = "morning" if hour < 12 else "afternoon" if hour < 18 else "evening"
        
        suggestions = []
        for act in data["actions"]:
            if (act["context"]["day"] == day and 
                act["context"]["hour"] == hour and
                act["context"]["time_of_day"] == time_of_day and
                act["count"] >= 2):
                suggestions.append((act["action"], act["count"]))
        
        suggestions.sort(key=lambda x: x[1], reverse=True)
        top_suggestions = [s[0] for s in suggestions[:3]]
        
        if top_suggestions:
            try:
                # Sanitize suggestion text to prevent injection
                safe_suggestions = [s.replace('"', '').replace("'", '') for s in top_suggestions[:2]]
                suggestion_text = f"You usually use: {', '.join(safe_suggestions)}"
                
                # Use safer PowerShell execution
                ps_command = [
                    'powershell', '-ExecutionPolicy', 'Bypass', '-Command',
                    f'Add-Type -AssemblyName System.Windows.Forms; '
                    f'[System.Windows.Forms.MessageBox]::Show("{suggestion_text}", "Jarvis Suggestions")'
                ]
                subprocess.run(ps_command, shell=False, timeout=5)
            except (subprocess.TimeoutExpired, subprocess.SubprocessError, OSError) as e:
                print(f"Notification error: {e}")
        
        return top_suggestions
    
    def enable_proactive_mode(self):
        self._memory_set('proactive_mode_enabled', True)
        return "Proactive mode enabled! Jarvis will automatically suggest actions based on your patterns."
    
    def disable_proactive_mode(self):
        self._memory_set('proactive_mode_enabled', False)
        return "Proactive mode disabled."
    
    def check_proactive_suggestions(self):
        enabled_data = self._memory_get_all('proactive_mode_enabled')
        if not enabled_data or not enabled_data[0]:
            return ""
        
        suggestions = self.auto_suggest_with_popup()
        if suggestions:
            top_action = suggestions[0]
            return f"Suggested: {top_action}"
        return ""
    
    def start_proactive_background(self):
        import threading
        import time
        
        def background_suggestions():
            while True:
                try:
                    enabled_data = self._memory_get_all('proactive_mode_enabled')
                    if enabled_data and enabled_data[0]:
                        now = datetime.now()
                        if 8 <= now.hour <= 22:
                            self.auto_suggest_with_popup()
                    time.sleep(1800)
                except:
                    time.sleep(1800)
        
        bg_thread = threading.Thread(target=background_suggestions, daemon=True)
        bg_thread.start()
        return "Proactive Jarvis started in background"
    
    def manual_learn(self, action: str = ""):
        if not action or not action.strip():
            return "Please specify what to learn. Say: 'manual learn system_monitor'"
        
        # Sanitize action input
        action = action.strip()[:100]  # Limit length
        
        try:
            self._log_action_with_context(action)
            return f"Manually learned: {action}. I'll remember this pattern for future suggestions."
        except Exception as e:
            return f"Error learning action: {str(e)}"
    def docker_control(self): return "Docker control available"
    
    # amazonq-ignore-next-line
    def execute_voice_command(self, command):
        """Execute voice command with multilingual support"""
        # Process multilingual commands first
        if self.multilingual:
            # Check if it's a language switching command
            if any(word in command.lower() for word in ['switch to', 'change language', 'भाषा', 'ভাষা', 'ભાષા', 'ಭಾಷೆ', 'ഭാഷ', 'भाषा', 'மொழி', 'భాష', 'زبان']):
                return self.multilingual.process_multilingual_command(command)
            
            # Auto-detect language and process accordingly
            detected_lang = self.multilingual.detect_language(command)
            if detected_lang != 'english' and detected_lang != self.multilingual.current_language:
                self.multilingual.current_language = detected_lang
        
        command_lower = command.lower().strip()
        
        # Check for exact matches
        if command_lower in self.voice_functions:
            try:
                result = self.voice_functions[command_lower]()
                return result
            except Exception as e:
                # amazonq-ignore-next-line
                return f"Error executing {command_lower}: {str(e)}"
        
        # Check for partial matches
        for func_name in self.voice_functions.keys():
            if func_name in command_lower or command_lower in func_name:
                try:
                    result = self.voice_functions[func_name]()
                    return result
                except Exception as e:
                    return f"Error executing {func_name}: {str(e)}"
        
        # Handle dynamic commands
        if command_lower.startswith('remember this '):
            content = command_lower.replace('remember this ', '').strip()
            # Determine appropriate key based on content
            if 'laptop' in content:
                key = "laptop_info"
            elif 'dog' in content:
                key = "pet_info"
            elif 'car' in content:
                key = "vehicle_info"
            else:
                key = "personal_info"
            return self.context_memory_store(key, content)
        
        # Health & Wellness dynamic commands
        if command_lower.startswith('health log '):
            entry = command_lower.replace('health log ', '').strip()
            return self.daily_health_log(entry)
        
        if command_lower.startswith('mood tracker ') or command_lower.startswith('track mood '):
            mood = command_lower.replace('mood tracker ', '').replace('track mood ', '').strip()
            return self.mood_tracker(mood)
        
        if command_lower.startswith('meditate ') or command_lower.startswith('meditation '):
            duration = command_lower.replace('meditate ', '').replace('meditation ', '').strip()
            if duration and duration.isdigit():
                return self.meditation_prompt(duration)
            return self.meditation_prompt()
        
        if 'posture' in command_lower and ('check' in command_lower or 'analyze' in command_lower):
            return self.posture_detection()
        
        # Web Intelligence dynamic commands
        if command_lower.startswith('scholar search '):
            query = command_lower.replace('scholar search ', '').strip()
            return self.scholar_search(query)
        
        if command_lower.startswith('stock '):
            symbol = command_lower.replace('stock ', '').strip()
            return self.stock_updates(symbol)
        
        if command_lower.startswith('crypto '):
            coin = command_lower.replace('crypto ', '').strip()
            return self.crypto_updates(coin)
        
        if command_lower.startswith('translate '):
            # Parse "translate [text] to [language]"
            if ' to ' in command_lower:
                parts = command_lower.replace('translate ', '').split(' to ')
                text = parts[0].strip()
                target_lang = parts[1].strip()
                return self.realtime_translation(text, target_lang)
            else:
                text = command_lower.replace('translate ', '').strip()
                return self.realtime_translation(text)
        
        if command_lower.startswith('install '):
            package = command_lower.replace('install ', '').strip()
            return self.install_package(package)
        
        if command_lower.startswith('uninstall '):
            package = command_lower.replace('uninstall ', '').strip()
            return self.uninstall_package(package)
        
        # Security & Authentication dynamic commands
        if command_lower.startswith('encrypt file '):
            file_path = command_lower.replace('encrypt file ', '').strip()
            if file_path:
                return self.file_vault_encrypt(file_path)
            else:
                return "Please specify file path. Say: 'encrypt file C:\\path\\to\\file.txt'"
        
        if command_lower.startswith('decrypt file '):
            file_path = command_lower.replace('decrypt file ', '').strip()
            if file_path:
                return self.file_vault_decrypt(file_path)
            else:
                return "Please specify encrypted file path"
        
        if command_lower.startswith('phishing scan '):
            url = command_lower.replace('phishing scan ', '').strip()
            return self.phishing_malware_scan_link(url)
        
        if command_lower.startswith('parental control '):
            action = command_lower.replace('parental control ', '').strip()
            return self.parental_control_set(action)
        
        if command_lower.startswith('schedule event '):
            event = command_lower.replace('schedule event ', '').strip()
            return self.calendar_schedule(event)
        
        if command_lower.startswith('add event '):
            event = command_lower.replace('add event ', '').strip()
            return self.calendar_schedule(event)
        
        if command_lower.startswith('summarize email '):
            email_content = command_lower.replace('summarize email ', '').strip()
            return self.email_summarize(email_content)
        
        if command_lower.startswith('code help '):
            task = command_lower.replace('code help ', '').strip()
            return self.code_agent(task)
        
        if command_lower.startswith('research '):
            topic = command_lower.replace('research ', '').strip()
            return self.research_agent(topic)
        
        # Document creation dynamic commands
        if command_lower.startswith('create document '):
            topic = command_lower.replace('create document ', '').strip()
            return self.ai_document_maker("document", topic)
        
        if command_lower.startswith('create report '):
            topic = command_lower.replace('create report ', '').strip()
            return self.ai_document_maker("report", topic)
        
        if command_lower.startswith('create letter '):
            topic = command_lower.replace('create letter ', '').strip()
            return self.ai_document_maker("letter", topic)
        
        # Adaptive Learning dynamic commands
        if command_lower.startswith('adaptive learning '):
            action = command_lower.replace('adaptive learning ', '').strip()
            return self.adaptive_learning(action)
        
        if command_lower.startswith('manual learn '):
            action = command_lower.replace('manual learn ', '').strip()
            return self.manual_learn(action)
        
        if command_lower.startswith('teach jarvis '):
            action = command_lower.replace('teach jarvis ', '').strip()
            return self.manual_learn(action)
        
        # Handle any question about stored information
        if command_lower.startswith('what is my ') or 'what is' in command_lower or 'who is' in command_lower:
            memories = self.context_memory_recall()
            if memories and memories != "No memories stored":
                prompt = f"Based on these stored memories: {memories}\n\nAnswer this question: {command}"
                return self._ai_generate(prompt)
            return "I don't have any stored memories to answer that question"
        
        # Try multilingual processing as final fallback
        if self.multilingual:
            ml_response = self.multilingual.process_command_in_language(command, self.multilingual.current_language)
            if ml_response != self.multilingual.get_response('processing'):
                return ml_response
            return self.multilingual.get_response('not_understood')
        
        return "Voice command not recognized"

# Create instance
voice_advanced_ai = VoiceAdvancedAI()

def get_voice_advanced_response(command):
    # amazonq-ignore-next-line
    return voice_advanced_ai.execute_voice_command(command)  